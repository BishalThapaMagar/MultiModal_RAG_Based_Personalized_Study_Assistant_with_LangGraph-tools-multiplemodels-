import json
import subprocess
from fastapi import FastAPI, Query, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches
from fastapi.responses import JSONResponse
from pathlib import Path
import uuid
import requests
import re
import random
from fastapi.middleware.cors import CORSMiddleware

# ----------------------------
# Setup FastAPI and directories
# ----------------------------
app = FastAPI(title="Dynamic AI PPT Generator", version="5.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # change to frontend URL in prod
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
OUTPUT_DIR = Path("outputs")
OUTPUT_DIR.mkdir(exist_ok=True)
IMAGE_DIR = Path("images")
IMAGE_DIR.mkdir(exist_ok=True)
TEMPLATE_DIR = Path("templates")
TEMPLATE_DIR.mkdir(exist_ok=True)

PLACEHOLDER_IMAGE = IMAGE_DIR / "placeholder.png"

# ----------------------------
# Topic -> Template mapping
# ----------------------------
TEMPLATE_MAPPING = {
    "business": "business_template.pptx",
    "finance": "business_template.pptx",
    "marketing": "business_template.pptx",
    "science": "science_template.pptx",
    "education": "education_template.pptx",
    "teaching": "education_template.pptx",
    "technology": "tech_template.pptx",
    "ai": "tech_template.pptx",
    "machine learning": "tech_template.pptx"
}


def choose_template(topic: str) -> Path:
    """
    Chooses a PPT template based on the topic.
    Falls back to default_template.pptx if no match is found.
    """
    topic_lower = topic.lower()
    for key, template_file in TEMPLATE_MAPPING.items():
        if key in topic_lower:
            template_path = TEMPLATE_DIR / template_file
            if template_path.exists():
                return template_path
            else:
                print(f"Template {template_file} for topic '{key}' not found, will try default.")

    # Fallback to default template
    default_template = TEMPLATE_DIR / "default_template.pptx"
    if default_template.exists():
        print("Using default template as no specific template was found.")
        return default_template
    else:
        # If default template also doesn't exist, raise an error
        raise FileNotFoundError("No valid PPT template found, including default_template.pptx!")


# ----------------------------
# Generate slide outline using Groq API
# ----------------------------
import sys
sys.path.append('../..')
from shared.GroqClient import generate_completion

def get_slide_outline(topic: str, num_slides: int = 5):
    prompt = f"""
    Create a {num_slides}-slide PowerPoint presentation outline on the topic: "{topic}".
    Return the result in valid JSON only in the format:
    {{
      "title": "Main Presentation Title",
      "slides": [
        {{"title": "Slide 1 Title", "content": "Bullet point 1. Bullet point 2."}},
        ...
      ]
    }}
    """
    
    text_output = generate_completion(prompt, "You are an AI assistant that creates presentation outlines. Always return valid JSON.")
    match = re.search(r'\{.*\}', text_output, re.DOTALL)
    if match:
        try:
            data = json.loads(match.group())
            return data
        except json.JSONDecodeError:
            print("Error decoding JSON from LLaMA output")
            print("Raw output:", text_output)
            return None
    return None

# ----------------------------
# Download topic-related image
# ----------------------------
def get_topic_image(topic: str) -> Path:
    search_url = f"https://source.unsplash.com/800x600/?{topic.replace(' ', ',')}"
    img_path = IMAGE_DIR / f"{uuid.uuid4().hex}.jpg"
    try:
        r = requests.get(search_url, allow_redirects=True, timeout=10)
        if r.status_code == 200 and r.headers['Content-Type'].startswith('image'):
            with open(img_path, "wb") as f:
                f.write(r.content)
            return img_path
    except Exception as e:
        print(f"Error downloading image for topic '{topic}': {e}")
    # fallback
    if PLACEHOLDER_IMAGE.exists():
        return PLACEHOLDER_IMAGE
    return None

# ----------------------------
# Helper: add content with proper bullets
# ----------------------------
def add_bullet_points(placeholder, content: str):
    placeholder.text = ""  # clear existing
    # split content into bullets (handles • or â€¢ or newlines)
    bullets = re.split(r'[•â€¢\n]+', content)
    for b in bullets:
        b = b.strip()
        if b:
            p = placeholder.text_frame.add_paragraph()
            p.text = b
            p.level = 0  # top-level bullet

# ----------------------------
# Create PPT with dynamic layouts
# ----------------------------
def create_ppt(slide_data, output_file: Path, template_file: Path):
    prs = Presentation(template_file)

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = slide_data["title"]
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "By Group 21"

    # Content Slides
    for s in slide_data["slides"]:
        layout_choice = random.choice(["text_only", "image_left", "image_right", "full_image"])
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = s["title"]

        # Add content with proper bullets
        if len(slide.placeholders) > 1:
            add_bullet_points(slide.placeholders[1], s["content"])

        # Add topic image according to layout
        img_path = get_topic_image(s["title"])
        if img_path and img_path.exists():
            if layout_choice == "image_left":
                slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(2), Inches(4), Inches(3))
            elif layout_choice == "image_right":
                slide.shapes.add_picture(str(img_path), Inches(5), Inches(2), Inches(4), Inches(3))
            elif layout_choice == "full_image":
                slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
            # text_only does not add image

    prs.save(output_file)
    return output_file

# ----------------------------
# FastAPI Endpoints
# ----------------------------
@app.get("/")
def home():
    return {"message": "Welcome to Dynamic AI PPT Generator with proper bullets 🚀"}

@app.get("/generate_ppt")
def generate_ppt(topic: str = Query(..., description="Topic of the PPT"),
                 slides: int = Query(5, description="Number of slides")):
    outline = get_slide_outline(topic, num_slides=slides)
    if not outline:
        raise HTTPException(status_code=500, detail="Failed to generate PPT outline")

    try:
        template_file = choose_template(topic)
    except FileNotFoundError as e:
        raise HTTPException(status_code=500, detail=str(e))

    safe_topic = re.sub(r'[^\w\d-]', '_', topic)
    output_file = OUTPUT_DIR / f"{uuid.uuid4().hex}_{safe_topic}.pptx"
    ppt_file = create_ppt(outline, output_file, template_file)

    return FileResponse(
        ppt_file,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=ppt_file.name
    )