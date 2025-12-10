import json
import uuid
from pathlib import Path
from pptx import Presentation

def get_llm_function(model_provider: str):
    """Route to appropriate AI service"""
    if model_provider == "groq":
        from ai_services.GroqClient import generate_completion
        return generate_completion
    elif model_provider == "groq_pro":
        from ai_services.GroqClient import generate_completion_paid
        return generate_completion_paid
    elif model_provider == "gemini":
        from ai_services.GeminiClient import call_gemini
        return call_gemini
    elif model_provider == "perplexity":
        from ai_services.PerplexityClient import call_perplexity_chat
        return call_perplexity_chat
    else:
        # Default to groq
        from ai_services.GroqClient import generate_completion
        return generate_completion

def get_slide_outline(topic: str, num_slides: int, model_provider: str = "groq") -> dict:
    """Generate slide outline using selected AI model"""
    llm_func = get_llm_function(model_provider)
    
    prompt = f"""Create a {num_slides}-slide presentation outline on: "{topic}"
    
Return ONLY valid JSON in this exact format:
{{
  "title": "Main Title Here",
  "slides": [
    {{"title": "Slide 1 Title", "content": "Point 1. Point 2. Point 3."}},
    {{"title": "Slide 2 Title", "content": "Point 1. Point 2. Point 3."}}
  ]
}}"""
    
    response = llm_func(prompt, "You are a presentation expert. Return only valid JSON, no markdown formatting.")
    
    # Clean the response
    response = response.strip()
    if response.startswith("```json"):
        response = response[7:]
    if response.startswith("```"):
        response = response[3:]
    if response.endswith("```"):
        response = response[:-3]
    response = response.strip()
    
    try:
        return json.loads(response)
    except:
        # Fallback outline if parsing fails
        return {
            "title": f"{topic} - Presentation",
            "slides": [
                {"title": "Introduction", "content": f"Overview of {topic}. Key concepts. Importance."},
                {"title": "Main Content 1", "content": "First key point. Supporting details. Examples."},
                {"title": "Main Content 2", "content": "Second key point. Supporting details. Examples."},
                {"title": "Applications", "content": "Real-world uses. Case studies. Benefits."},
                {"title": "Conclusion", "content": "Summary of key points. Future implications. Call to action."}
            ][:num_slides]
        }

def generate_ppt(topic: str, num_slides: int = 5, model_provider: str = "groq") -> dict:
    """Generate PowerPoint presentation"""
    current_dir = Path(__file__).parent.parent
    
    # Generate outline
    outline = get_slide_outline(topic, num_slides, model_provider)
    
    if not outline:
        raise Exception("Failed to generate outline")
    
    # Handle template selection
    template_dir = current_dir / "ppt-generator-" / "backend" / "templates"
    template_mapping = {
        "ai": "tech_template.pptx",
        "technology": "tech_template.pptx", 
        "tech": "tech_template.pptx",
        "business": "business_template.pptx",
        "science": "science_template.pptx"
    }
    
    topic_lower = topic.lower()
    template_name = "default_template.pptx"
    
    for key, tmpl in template_mapping.items():
        if key in topic_lower:
            template_name = tmpl
            break
    
    template_file = template_dir / template_name
    if not template_file.exists():
        template_file = template_dir / "default_template.pptx"
    
    if not template_file.exists():
        raise Exception(f"Template not found: {template_file}")
    
    # Create output directory
    output_dir = current_dir / "outputs"
    output_dir.mkdir(exist_ok=True)
    
    safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).replace(' ', '_')
    output_filename = f"{safe_topic}_{uuid.uuid4().hex[:8]}.pptx"
    output_file = output_dir / output_filename
    
    # Create PPT
    if template_file.exists():
        prs = Presentation(template_file)
    else:
        prs = Presentation()
    
    # Clear existing slides except first
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1]
    
    # Update title slide
    if len(prs.slides) > 0:
        title_slide = prs.slides[0]
        if title_slide.shapes.title:
            title_slide.shapes.title.text = outline.get("title", topic)
    else:
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = outline.get("title", topic)
    
    # Add content slides
    for slide_data in outline.get("slides", []):
        content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data.get("title", "")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_data.get("content", "")
    
    # Save presentation
    prs.save(output_file)
    
    return {
        "success": True,
        "filename": output_filename,
        "download_url": f"/download/{output_filename}",
        "slides_count": len(outline.get("slides", [])),
        "title": outline.get("title", topic),
        "model_used": model_provider
    }

